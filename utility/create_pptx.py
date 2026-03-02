# Create a PowerPoint diagram based on the provided structure

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_CONNECTOR

# Create presentation
prs = Presentation()
slide_layout = prs.slide_layouts[6]  # Blank layout
slide = prs.slides.add_slide(slide_layout)

# Helper function to add a box
def add_box(slide, text, left, top, width=1.8, height=0.8):
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    box.text_frame.text = text
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    box.fill.background()  # No fill
    box.line.color.rgb = RGBColor(0, 0, 0)
    return box

# Add nodes
leader = add_box(slide, "Leader: Dan", 4, 0.5)

cumhur = add_box(slide, "Cumhur", 1.5, 2)
assistant = add_box(slide, "Missing Assistant #1", 4, 2)
sofia = add_box(slide, "Sofia", 6.5, 2)

razvan = add_box(slide, "Razvan", 1, 3.5)
anders = add_box(slide, "Anders", 2.5, 3.5)
peter = add_box(slide, "Peter", 4, 3.5)
prithvi = add_box(slide, "Prithvi", 7, 3.5)

# Helper function to connect shapes
def connect_shapes(slide, shape1, shape2):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        shape1.left + shape1.width / 2,
        shape1.top + shape1.height,
        shape2.left + shape2.width / 2,
        shape2.top
    )
    connector.line.color.rgb = RGBColor(0, 0, 0)

# Connections
connect_shapes(slide, leader, cumhur)
connect_shapes(slide, leader, assistant)
connect_shapes(slide, leader, sofia)

connect_shapes(slide, cumhur, razvan)
connect_shapes(slide, cumhur, anders)

connect_shapes(slide, assistant, peter)
connect_shapes(slide, sofia, prithvi)

# Save file
file_path = "/Users/cer/Downloads/Team_Structure_Diagram.pptx"
prs.save(file_path)

file_path